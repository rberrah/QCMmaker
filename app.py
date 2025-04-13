import streamlit as st
import os
import json
from datetime import datetime
from transformers import pipeline
import pandas as pd

# Modules pour extraire du texte depuis différents formats
import PyPDF2
import docx
from pptx import Presentation

# =============================================
# Définition des fichiers de stockage persistants
COURSES_DB_FILE = "courses_db.json"
EXAMS_DB_FILE = "exams_db.json"
FEEDBACK_DB_FILE = "feedback_history.json"

# =============================================
# Fonctions de stockage pour les cours
def load_courses_db():
    if os.path.exists(COURSES_DB_FILE):
        try:
            with open(COURSES_DB_FILE, "r") as f:
                courses = json.load(f)
        except Exception as e:
            st.error("Erreur lors du chargement de la base de cours.")
            courses = []
    else:
        courses = []
    return courses

def save_courses_db(courses):
    with open(COURSES_DB_FILE, "w") as f:
        json.dump(courses, f, ensure_ascii=False, indent=4)

# =============================================
# Fonctions de stockage pour les annales
def load_exams_db():
    if os.path.exists(EXAMS_DB_FILE):
        try:
            with open(EXAMS_DB_FILE, "r") as f:
                exams = json.load(f)
        except Exception as e:
            st.error("Erreur lors du chargement de la base de sujets d'annales.")
            exams = []
    else:
        exams = []
    return exams

def save_exams_db(exams):
    with open(EXAMS_DB_FILE, "w") as f:
        json.dump(exams, f, ensure_ascii=False, indent=4)

# =============================================
# Fonctions de stockage pour le feedback
def load_feedback_history():
    if os.path.exists(FEEDBACK_DB_FILE):
        try:
            with open(FEEDBACK_DB_FILE, "r") as f:
                history = json.load(f)
        except Exception as e:
            st.error("Erreur lors du chargement de l'historique des feedbacks.")
            history = []
    else:
        history = []
    return history

def save_feedback_entry(entry):
    history = load_feedback_history()
    history.append(entry)
    with open(FEEDBACK_DB_FILE, "w") as f:
        json.dump(history, f, ensure_ascii=False, indent=4)

# =============================================
# Fonction d'extraction de texte depuis différents formats
def extract_text_from_file(uploaded_file):
    """
    Extrait le texte d'un fichier en fonction de son extension.
    Formats supportés : .txt, .pdf, .docx, .pptx, .xls et .xlsx.
    Pour les PDF, chaque page est précédée d'une annotation indiquant
    le numéro de page et le nom du fichier.
    """
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    text = ""
    if file_extension == ".txt":
        try:
            text = uploaded_file.getvalue().decode("utf-8")
        except Exception as e:
            text = "Erreur lors de la lecture du fichier texte."
    elif file_extension == ".pdf":
        try:
            reader = PyPDF2.PdfReader(uploaded_file)
            for idx, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text += f"Page {idx+1} de '{uploaded_file.name}':\n" + page_text + "\n\n"
        except Exception as e:
            text = "Erreur lors de l'extraction du PDF."
    elif file_extension == ".docx":
        try:
            doc = docx.Document(uploaded_file)
            text = "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            text = "Erreur lors de l'extraction du document Word."
    elif file_extension == ".pptx":
        try:
            presentation = Presentation(uploaded_file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            text = "Erreur lors de l'extraction du PowerPoint."
    elif file_extension in [".xls", ".xlsx"]:
        try:
            df = pd.read_excel(uploaded_file)
            text = df.to_string(index=False)
        except Exception as e:
            text = "Erreur lors de l'extraction du fichier Excel."
    else:
        text = "Format de fichier non supporté."
    return text

# =============================================
# Initialisation du générateur de questions (ici modèle T5-base)
generator = pipeline("text2text-generation", model="t5-base")

# =============================================
# Interface Streamlit avec plusieurs onglets
st.title("Générateur de Questions Type Annales")
tabs = st.tabs(["Cours", "Sujets d'annales", "Générer les questions", "Feedback", "Historique Feedback"])

# ---------- Onglet 1 : Déposer un cours ----------
with tabs[0]:
    st.header("Déposer un cours")
    mode_course = st.radio("Mode de dépôt :", ("Téléverser des fichiers", "Saisie manuelle"), key="mode_course")
    course_year = st.number_input("Année du cours", value=2025, step=1, key="course_year")
    
    if mode_course == "Téléverser des fichiers":
        common_title_course = st.text_input("Titre commun du cours (optionnel)", key="common_course_title")
        uploaded_course_files = st.file_uploader("Téléversez un ou plusieurs fichiers de cours",
                                                 type=["txt", "pdf", "docx", "pptx", "xls", "xlsx"],
                                                 accept_multiple_files=True,
                                                 key="uploaded_course_files")
        if uploaded_course_files:
            new_courses = []
            if common_title_course:
                combined_text = ""
                for f in uploaded_course_files:
                    text = extract_text_from_file(f)
                    if text.startswith("Erreur") or text.startswith("Format"):
                        st.error(f"{f.name} : {text}")
                    else:
                        combined_text += text + "\n"
                if combined_text:
                    new_courses.append({
                        "title": common_title_course,
                        "text": combined_text,
                        "year": int(course_year)
                    })
            else:
                for f in uploaded_course_files:
                    text = extract_text_from_file(f)
                    if text.startswith("Erreur") or text.startswith("Format"):
                        st.error(f"{f.name} : {text}")
                    else:
                        new_courses.append({
                            "title": f.name,
                            "text": text,
                            "year": int(course_year)
                        })
            if new_courses:
                courses_db = load_courses_db()
                courses_db.extend(new_courses)
                save_courses_db(courses_db)
                st.success(f"{len(new_courses)} cours enregistrés avec succès.")
    else:
        manual_course_title = st.text_input("Titre du cours ou chapitre", key="manual_course_title")
        manual_course_text = st.text_area("Collez ici le contenu du cours", height=300, key="manual_course_text")
        if manual_course_text:
            new_course = {
                "title": manual_course_title if manual_course_title else "Cours manuel",
                "text": manual_course_text,
                "year": int(course_year)
            }
            courses_db = load_courses_db()
            courses_db.append(new_course)
            save_courses_db(courses_db)
            st.success("Cours saisi manuellement enregistré.")

# ---------- Onglet 2 : Déposer un sujet d'annales ----------
with tabs[1]:
    st.header("Déposer un sujet d'annale")
    mode_exam = st.radio("Mode de dépôt :", ("Téléverser des fichiers", "Saisie manuelle"), key="mode_exam")
    exam_year = st.number_input("Année du sujet d'annale", value=2025, step=1, key="exam_year")
    
    if mode_exam == "Téléverser des fichiers":
        common_title_exam = st.text_input("Titre commun du sujet d'annale (optionnel)", key="common_exam_title")
        uploaded_exam_files = st.file_uploader("Téléversez un ou plusieurs fichiers de sujets d'annales",
                                               type=["txt", "pdf", "docx", "pptx", "xls", "xlsx"],
                                               accept_multiple_files=True,
                                               key="uploaded_exam_files")
        if uploaded_exam_files:
            new_exams = []
            if common_title_exam:
                combined_text = ""
                for f in uploaded_exam_files:
                    text = extract_text_from_file(f)
                    if text.startswith("Erreur") or text.startswith("Format"):
                        st.error(f"{f.name} : {text}")
                    else:
                        combined_text += text + "\n"
                if combined_text:
                    new_exams.append({
                        "title": common_title_exam,
                        "text": combined_text,
                        "year": int(exam_year)
                    })
            else:
                for f in uploaded_exam_files:
                    text = extract_text_from_file(f)
                    if text.startswith("Erreur") or text.startswith("Format"):
                        st.error(f"{f.name} : {text}")
                    else:
                        new_exams.append({
                            "title": f.name,
                            "text": text,
                            "year": int(exam_year)
                        })
            if new_exams:
                exams_db = load_exams_db()
                exams_db.extend(new_exams)
                save_exams_db(exams_db)
                st.success(f"{len(new_exams)} sujets d'annales enregistrés avec succès.")
    else:
        manual_exam_title = st.text_input("Titre du sujet d'annale", key="manual_exam_title")
        manual_exam_text = st.text_area("Collez ici le contenu du sujet d'annale", height=200, key="manual_exam_text")
        if manual_exam_text:
            new_exam = {
                "title": manual_exam_title if manual_exam_title else "Sujet manuel",
                "text": manual_exam_text,
                "year": int(exam_year)
            }
            exams_db = load_exams_db()
            exams_db.append(new_exam)
            save_exams_db(exams_db)
            st.success("Sujet d'annale saisi manuellement enregistré.")

# ---------- Onglet 3 : Générer les questions ----------
with tabs[2]:
    st.header("Générer les questions")
    
    st.subheader("Filtrer par année")
    start_year = st.number_input("Année de début", value=2000, step=1, key="start_year")
    end_year = st.number_input("Année de fin", value=2025, step=1, key="end_year")
    
    # Filtrer les cours selon l'année
    courses_db = load_courses_db()
    filtered_courses = [course for course in courses_db if start_year <= course["year"] <= end_year]
    course_options = {}
    for course in filtered_courses:
        display_title = f"{course['title']} ({course['year']})"
        course_options[display_title] = course["text"]
    selected_courses = st.multiselect("Sélectionnez les cours à utiliser", options=list(course_options.keys()), key="selected_courses_gen")
    course_content = ""
    for course in selected_courses:
        course_content += f"Cours : {course}\n{course_options[course]}\n\n"
    
    # Filtrer les sujets d'annales selon l'année
    exams_db = load_exams_db()
    filtered_exams = [exam for exam in exams_db if start_year <= exam["year"] <= end_year]
    exam_options = {}
    for exam in filtered_exams:
        display_title = f"{exam['title']} ({exam['year']})"
        exam_options[display_title] = exam["text"]
    selected_exams = st.multiselect("Sélectionnez les sujets d'annales à utiliser", options=list(exam_options.keys()), key="selected_exams_gen")
    exam_content = ""
    for exam in selected_exams:
        exam_content += f"Sujet d'annale : {exam}\n{exam_options[exam]}\n\n"
    
    if not course_content or not exam_content:
        st.info("Veuillez sélectionner au moins un cours et un sujet d'annale dans la plage d'années spécifiée.")
    else:
        generation_mode = st.radio("Choisissez le mode de génération :", ["Une Question", "Annales Complète"], key="generation_mode")
        base_prompt = "Génère des questions d'examen basées sur le contenu suivant :\n\n"
        base_prompt += f"{course_content}\n"
        base_prompt += f"{exam_content}\n"
        if generation_mode == "Une Question":
            base_prompt += "\nVeuillez générer une seule question d'examen. Pour cette question, indiquez la réponse ainsi que la provenance sous la forme 'Page X de \"Nom du document\"'."
            max_len = 128
        else:
            base_prompt += "\nVeuillez générer une annale complète composée de plusieurs questions d'examen. Pour chacune, indiquez la réponse ainsi que la provenance sous la forme 'Page X de \"Nom du document\"'."
            max_len = 256
        
        st.write("Prompt généré automatiquement (modifiable) :")
        st.text_area("Prompt (modifiable)", value=base_prompt, key="custom_prompt", height=150)
        
        if st.button("Générer"):
            prompt_to_use = st.session_state.custom_prompt
            st.write("Génération en cours...")
            with st.spinner("Veuillez patienter..."):
                try:
                    output = generator(prompt_to_use, max_length=max_len, num_return_sequences=1)
                    generated_questions = output[0]["generated_text"]
                    st.subheader("Questions générées")
                    st.write(generated_questions)
                    st.session_state.generated_questions = generated_questions
                    # Sauvegarde des titres utilisés (sans l'année)
                    used_courses = [c.split(" (")[0] for c in selected_courses]
                    st.session_state.generated_questions_courses = used_courses
                    if not selected_courses and course_options:
                        st.session_state.generated_questions_courses = []
                except Exception as e:
                    st.error(f"Une erreur est survenue lors de la génération : {e}")

# ---------- Onglet 4 : Feedback ----------
with tabs[3]:
    st.header("Donner votre avis sur la qualité des questions")
    if "generated_questions" not in st.session_state:
        st.info("Veuillez d'abord générer des questions dans l'onglet 'Générer les questions'.")
    elif ("generated_questions_courses" not in st.session_state or 
          len(set(st.session_state.generated_questions_courses)) != 1):
        st.info("Le feedback est disponible uniquement lorsque les questions sont générées à partir d'un unique cours. Veuillez sélectionner un seul intitulé de cours lors de la génération.")
    else:
        st.write("Notez chaque question individuellement :")
        questions_list = [q.strip() for q in st.session_state.generated_questions.split("\n") if q.strip() != ""]
        question_ratings = {}
        for i, question in enumerate(questions_list):
            st.markdown(f"**Question {i+1}** : {question}")
            rating = st.slider(
                f"Votre note pour la question {i+1} (1 = Médiocre, 5 = Excellent)",
                min_value=1,
                max_value=5,
                value=3,
                step=1,
                key=f"rating_{i}"
            )
            question_ratings[f"Question {i+1}"] = {"question": question, "rating": rating}
        overall_rating = st.slider("Note globale pour l'ensemble des questions (1 = Médiocre, 5 = Excellent)",
                                   min_value=1, max_value=5, value=3, step=1, key="overall_rating")
        if st.button("Envoyer votre feedback"):
            feedback_entry = {
                "timestamp": datetime.now().isoformat(),
                "generated_questions": st.session_state.generated_questions,
                "question_ratings": question_ratings,
                "overall_rating": overall_rating,
                "courses_used": st.session_state.generated_questions_courses
            }
            save_feedback_entry(feedback_entry)
            st.success("Merci pour votre feedback !")
            st.session_state.last_feedback = feedback_entry

# ---------- Onglet 5 : Historique des feedbacks ----------
with tabs[4]:
    st.header("Historique des feedbacks")
    history = load_feedback_history()
    if not history:
        st.info("Aucun feedback enregistré pour le moment.")
    else:
        all_courses = set()
        for entry in history:
            for course_title in entry.get("courses_used", []):
                all_courses.add(course_title)
        all_courses = list(all_courses)
        all_courses.sort()
        selected_filter = st.selectbox("Filtrer par intitulé de cours", options=["Tous"] + all_courses, key="filter_course")
        for entry in history:
            entry_courses = entry.get("courses_used", [])
            if selected_filter != "Tous" and selected_filter not in entry_courses:
                continue
            st.markdown(f"**Date et heure :** {entry.get('timestamp', 'Inconnue')}")
            st.markdown(f"**Intitulé(s) du cours utilisé(s) :** {', '.join(entry_courses) if entry_courses else 'N/A'}")
            st.markdown(f"**Note Globale :** {entry.get('overall_rating', 'N/A')}")
            qr = entry.get("question_ratings", {})
            if qr:
                st.markdown("**Notes par question :**")
                for key, val in qr.items():
                    st.markdown(f"- {key} : {val.get('rating', 'N/A')} (Question : {val.get('question', '')})")
            st.markdown("---")
